from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_presentation():
    prs = Presentation()

    # 主题颜色
    PRIMARY_COLOR = RGBColor(0, 224, 208)
    SECONDARY_COLOR = RGBColor(0, 160, 144)
    ACCENT_COLOR = RGBColor(0, 100, 255)
    BG_DARK = RGBColor(10, 10, 10)
    BG_LIGHT = RGBColor(26, 26, 26)
    TEXT_WHITE = RGBColor(255, 255, 255)
    TEXT_GRAY = RGBColor(180, 180, 180)

    def set_slide_background(slide, color=BG_DARK):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title, subtitle="", author=""):
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, BG_DARK)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True

        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        for p in subtitle_shape.text_frame.paragraphs:
            p.font.color.rgb = TEXT_GRAY
            p.font.size = Pt(18)

        if author:
            txBox = slide.shapes.add_textbox(Inches(9), Inches(6), Inches(3), Inches(0.5))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = author
            p.font.color.rgb = TEXT_GRAY
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.RIGHT

        return slide

    def add_content_slide(title, content_items=None, image_path=None, layout="title_content"):
        if layout == "title_content":
            slide_layout = prs.slide_layouts[1]
        elif layout == "title_only":
            slide_layout = prs.slide_layouts[5]
        else:
            slide_layout = prs.slide_layouts[1]

        slide = prs.slides.add_slide(slide_layout)
        set_slide_background(slide, BG_DARK)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        if content_items and layout != "title_only":
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()
            for i, item in enumerate(content_items):
                p = tf.add_paragraph()
                p.text = item
                if i == 0:
                    p.font.color.rgb = PRIMARY_COLOR
                    p.font.size = Pt(18)
                    p.font.bold = True
                else:
                    p.font.color.rgb = TEXT_GRAY
                    p.font.size = Pt(16)
                p.level = 0

        if image_path:
            slide.shapes.add_picture(image_path, Inches(5), Inches(1.5), width=Inches(4.5))

        return slide

    def add_bullet_slide(title, bullet_points):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        set_slide_background(slide, BG_DARK)

        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        content_shape = slide.placeholders[1]
        tf = content_shape.text_frame
        tf.clear()
        for i, item in enumerate(bullet_points):
            if len(item) == 4:
                level, text, is_bold, color = item
            elif len(item) == 3:
                level, text, is_bold = item
                color = None
            else:
                level, text = item
                is_bold = False
                color = None

            p = tf.add_paragraph()
            p.text = text
            p.level = level
            p.font.size = Pt(16)
            if level == 0:
                p.font.color.rgb = color if color else PRIMARY_COLOR
                p.font.bold = is_bold
            else:
                p.font.color.rgb = TEXT_GRAY
                p.font.bold = is_bold

        return slide

    def add_two_column_slide(title, left_content, right_content):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_background(slide, BG_DARK)

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_shape.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.color.rgb = PRIMARY_COLOR
        p.font.size = Pt(32)
        p.font.bold = True

        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(5))
        tf = left_box.text_frame
        for i, item in enumerate(left_content):
            p = tf.add_paragraph()
            p.text = item
            if i == 0:
                p.font.color.rgb = ACCENT_COLOR
                p.font.size = Pt(18)
                p.font.bold = True
            else:
                p.font.color.rgb = TEXT_GRAY
                p.font.size = Pt(15)

        right_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.2), Inches(5))
        tf = right_box.text_frame
        for i, item in enumerate(right_content):
            p = tf.add_paragraph()
            p.text = item
            if i == 0:
                p.font.color.rgb = ACCENT_COLOR
                p.font.size = Pt(18)
                p.font.bold = True
            else:
                p.font.color.rgb = TEXT_GRAY
                p.font.size = Pt(15)

        return slide

    def add_diagram_slide(title, diagram_items, colors=None):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_background(slide, BG_DARK)

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_shape.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.color.rgb = PRIMARY_COLOR
        p.font.size = Pt(32)
        p.font.bold = True

        box_width = Inches(2.8)
        box_height = Inches(1.5)
        start_x = Inches(0.6)
        start_y = Inches(1.5)
        gap_x = Inches(0.4)
        gap_y = Inches(0.3)

        if colors is None:
            colors = [PRIMARY_COLOR, SECONDARY_COLOR, ACCENT_COLOR, RGBColor(255, 165, 0), RGBColor(0, 255, 100)]

        for i, (label, desc) in enumerate(diagram_items):
            row = i // 3
            col = i % 3
            x = start_x + col * (box_width + gap_x)
            y = start_y + row * (box_height + gap_y)

            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = BG_LIGHT
            shape.line.color.rgb = colors[i % len(colors)]
            shape.line.width = Emu(300)

            txBox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.1), box_width - Inches(0.2), box_height - Inches(0.2))
            tf = txBox.text_frame
            p1 = tf.add_paragraph()
            p1.text = label
            p1.font.color.rgb = colors[i % len(colors)]
            p1.font.size = Pt(14)
            p1.font.bold = True

            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.color.rgb = TEXT_GRAY
            p2.font.size = Pt(11)

        return slide

    def add_code_slide(title, code_lines):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        set_slide_background(slide, BG_DARK)

        title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        tf = title_shape.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.color.rgb = PRIMARY_COLOR
        p.font.size = Pt(32)
        p.font.bold = True

        code_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5))
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(30, 30, 30)
        code_box.line.color.rgb = SECONDARY_COLOR

        tf = code_box.text_frame
        for line in code_lines:
            p = tf.add_paragraph()
            p.text = line
            p.font.color.rgb = TEXT_GRAY
            p.font.size = Pt(12)
            p.font.name = "Consolas"

        return slide

    # ========== 1. 封面页 ==========
    add_title_slide(
        "🚗 智能车载视觉系统",
        "基于 YOLOv8 的车道检测与距离估计系统",
        "Advanced Driver Assistance System"
    )

    # ========== 2. 项目概述 ==========
    add_bullet_slide(
        "项目概述",
        [
            (0, "项目简介", True),
            (1, "本系统是一个基于YOLOv8深度学习模型的智能车载视觉系统，集成了目标检测、车道线检测、距离估计和3D可视化四大核心模块", False),
            (0, "核心目标", True),
            (1, "实时识别前方道路环境中的车辆、行人、交通标志等目标", False),
            (1, "精准检测车道线并计算车道参数", False),
            (1, "基于单目相机实现目标距离估计与风险等级评估", False),
            (1, "通过3D可视化直观展示检测结果", False),
            (0, "应用价值", True),
            (1, "为高级驾驶辅助系统(ADAS)提供感知能力支撑", False),
            (1, "实现前方碰撞预警、车道偏离警告等安全功能", False),
        ]
    )

    # ========== 3. 系统架构 ==========
    add_diagram_slide(
        "系统架构",
        [
            ("输入模块", "摄像头/视频文件"),
            ("目标检测", "YOLOv8"),
            ("车道线检测", "OpenCV"),
            ("距离估计", "几何模型"),
            ("3D可视化", "立方体投影"),
            ("界面展示", "PyQt5"),
        ],
        colors=[
            PRIMARY_COLOR,
            RGBColor(0, 200, 255),
            RGBColor(0, 255, 100),
            RGBColor(255, 165, 0),
            RGBColor(255, 100, 200),
            RGBColor(100, 200, 255),
        ]
    )

    # ========== 4. 核心模块详解 ==========

    # 4.1 目标检测模块
    add_bullet_slide(
        "核心模块一：YOLOv8目标检测",
        [
            (0, "技术选型", True),
            (1, "采用 YOLOv8n 轻量化模型，平衡检测精度与推理速度", False),
            (1, "支持 COCO 数据集 80 个类别检测", False),
            (0, "检测目标类别", True),
            (1, "🚗 车辆类：car、truck、bus、motorcycle", False),
            (1, "🚶 行人类：person", False),
            (1, "🚦 交通类：traffic light、stop sign、fire hydrant", False),
            (0, "关键参数", True),
            (1, "置信度阈值：0.4（过滤低置信度检测）", False),
            (1, "IOU阈值：0.45（NMS非极大值抑制）", False),
            (0, "输出信息", True),
            (1, "类别ID、名称、置信度、边界框坐标、真实尺寸参数", False),
        ]
    )

    # 4.2 车道线检测模块
    add_bullet_slide(
        "核心模块二：车道线检测",
        [
            (0, "技术方案", True),
            (1, "传统计算机视觉方法：颜色阈值 + 边缘检测 + 滑动窗口 + 多项式拟合", False),
            (0, "处理流程", True),
            (1, "1. 颜色阈值分割：提取白色和黄色车道线像素", False),
            (1, "2. 边缘检测：Canny算子提取边缘特征", False),
            (1, "3. ROI掩码：限定前方道路区域，排除干扰", False),
            (1, "4. 透视变换：转换为鸟瞰图视角", False),
            (1, "5. 滑动窗口：搜索左右车道线像素位置", False),
            (1, "6. 多项式拟合：二次多项式描述车道曲线", False),
            (0, "输出信息", True),
            (1, "左右车道线坐标序列、中心线、曲率、车道宽度（米）", False),
        ]
    )

    # 4.3 距离估计模块
    add_bullet_slide(
        "核心模块三：单目距离估计",
        [
            (0, "技术原理", True),
            (1, "基于针孔相机模型和相似三角形原理", False),
            (1, "距离公式：z = (真实高度 × 焦距) / 像素高度", False),
            (0, "关键参数", True),
            (1, "相机焦距：根据视场角估算（70度 → 951像素）", False),
            (1, "目标真实尺寸：car(1.5m)、person(1.7m)、truck(3.5m)", False),
            (0, "优化策略", True),
            (1, "指数移动平均(EMA)：平滑帧间波动", False),
            (1, "分类处理：车辆/行人用高度，交通标志用宽度", False),
            (0, "风险等级分类", True),
            (1, "🔴 HIGH：<5米（危险）", False),
            (1, "🟠 MEDIUM：5-15米（警告）", False),
            (1, "🟡 LOW：15-30米（注意）", False),
            (1, "🟢 SAFE：>30米（安全）", False),
        ]
    )

    # 4.4 3D可视化模块
    add_bullet_slide(
        "核心模块四：3D可视化",
        [
            (0, "可视化策略", True),
            (1, "黑底风格，突出目标检测效果", False),
            (1, "不同目标类型采用差异化展示方式", False),
            (0, "目标可视化", True),
            (1, "🚗 车辆：12边立方体投影 + 地面投射点", False),
            (1, "🚶 行人：人体框 + 头部/脚部标记点", False),
            (1, "🚦 交通标志：矩形框 + 支撑柱", False),
            (0, "视图类型", True),
            (1, "原始视图：纯摄像头画面", False),
            (1, "车道线视图：原始+鸟瞰图", False),
            (1, "3D目标视图：黑底+3D立方体", False),
            (1, "综合视图：车道+3D+信息面板", False),
            (1, "信息面板：统计数据展示", False),
        ]
    )

    # ========== 5. 技术栈 ==========
    add_diagram_slide(
        "技术栈",
        [
            ("Python", "3.10+"),
            ("PyTorch", "深度学习框架"),
            ("Ultralytics", "YOLOv8实现"),
            ("OpenCV", "图像处理"),
            ("PyQt5", "图形界面"),
            ("NumPy/SciPy", "数值计算"),
        ],
        colors=[
            RGBColor(50, 100, 150),
            RGBColor(255, 100, 100),
            RGBColor(0, 200, 255),
            RGBColor(0, 200, 100),
            RGBColor(100, 100, 255),
            RGBColor(150, 150, 150),
        ]
    )

    # ========== 6. 项目结构 ==========
    add_content_slide(
        "项目结构",
        [
            "main.py          - 主程序入口，协调各模块",
            "detector.py      - YOLOv8目标检测模块",
            "lane_detector.py - 车道线检测模块",
            "distance_estimator.py - 距离估计模块",
            "visualizer.py    - 3D可视化模块",
            "ui.py            - PyQt5图形界面",
            "yolov8n.pt       - YOLOv8n模型权重",
            "requirements.txt - 依赖清单",
        ],
        layout="title_only"
    )

    # ========== 7. 功能演示 ==========
    add_two_column_slide(
        "功能演示 - 五大视图",
        [
            "📷 原始视图",
            "- 纯摄像头实时画面",
            "- 无任何叠加处理",
            "- 作为参考基准",
            "",
            "🛣️ 车道线检测视图",
            "- 上半：原始画面+车道线标注",
            "- 下半：鸟瞰图视角",
            "- 显示左右车道线和中心线",
        ],
        [
            "🎯 3D目标检测视图",
            "- 黑色背景",
            "- 3D立方体投影",
            "- 目标距离标注",
            "- 风险等级颜色标记",
            "",
            "🔗 综合视图",
            "- 融合车道线和目标检测",
            "- 3D立方体展示",
            "- 左侧信息面板",
            "- 完整道路环境感知",
        ]
    )

    # ========== 8. 运行方式 ==========
    add_content_slide(
        "快速开始",
        [
            "安装依赖：",
            "  pip install -r requirements.txt",
            "",
            "运行程序：",
            "  python main.py",
            "",
            "功能特性：",
            "  ✅ 支持摄像头实时采集",
            "  ✅ 支持视频文件播放",
            "  ✅ 置信度阈值可调",
            "  ✅ 截图保存功能",
            "  ✅ 五大视图独立展示",
        ],
        layout="title_only"
    )

    # ========== 9. 总结与展望 ==========
    add_bullet_slide(
        "总结与展望",
        [
            (0, "项目成果", True),
            (1, "成功实现基于YOLOv8的智能车载视觉系统原型", False),
            (1, "集成目标检测、车道线检测、距离估计、3D可视化", False),
            (1, "提供友好的PyQt5图形界面，支持多视图展示", False),
            (0, "未来改进方向", True),
            (1, "引入更先进的车道线检测算法（深度学习）", False),
            (1, "实现目标跟踪功能，提升距离估计稳定性", False),
            (1, "支持多摄像头融合", False),
            (1, "部署到嵌入式平台（如NVIDIA Jetson）", False),
            (0, "应用前景", True),
            (1, "自动驾驶领域的感知系统", False),
            (1, "智能交通监控系统", False),
            (1, "辅助驾驶教学系统", False),
        ]
    )

    # ========== 10. 致谢 ==========
    add_title_slide(
        "🎉 谢谢观看",
        "智能车载视觉系统 - 基于YOLOv8的车道检测与距离估计",
        "联系方式"
    )

    prs.save("/workspace/智能车载视觉系统演示文稿.pptx")
    print("✅ PowerPoint演示文稿已生成：智能车载视觉系统演示文稿.pptx")


if __name__ == "__main__":
    create_presentation()